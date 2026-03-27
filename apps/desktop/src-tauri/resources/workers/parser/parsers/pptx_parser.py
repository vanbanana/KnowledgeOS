from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .common import build_manifest, section, title_from_path


def parse_pptx(path: Path, source_type: str) -> dict:
    presentation = Presentation(str(path))
    slides_md = []
    sections = []

    for index, slide in enumerate(presentation.slides):
        texts = []
        slide_title = None
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if slide_title is None:
                slide_title = text
            texts.append(text)
        sections.append(section(f"slide-{index + 1}", index, slide_title or f"第 {index + 1} 页"))
        body = "\n\n".join(texts) if texts else "_本页未提取到文本_"
        slides_md.append(f"## 第 {index + 1} 页\n\n{body}")

    title = title_from_path(path)
    markdown = f"# {title}\n\n" + "\n\n".join(slides_md)
    return {
        "ok": True,
        "markdown": markdown.strip(),
        "manifest": build_manifest(
            title=title,
            source_type=source_type,
            source_path=str(path),
            sections=sections,
        ),
    }

