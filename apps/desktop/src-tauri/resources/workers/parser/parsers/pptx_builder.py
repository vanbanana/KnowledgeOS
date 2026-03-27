from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


PAGE_WIDTH = Inches(13.333)
PAGE_HEIGHT = Inches(7.5)
COLOR_BG = RGBColor(16, 22, 33)
COLOR_BG_SOFT = RGBColor(24, 33, 48)
COLOR_ACCENT = RGBColor(93, 148, 255)
COLOR_ACCENT_SOFT = RGBColor(42, 84, 160)
COLOR_TEXT = RGBColor(244, 247, 251)
COLOR_MUTED = RGBColor(176, 188, 206)
COLOR_PANEL = RGBColor(28, 38, 54)


def generate_pptx(output_path: str, presentation_json: str) -> dict[str, Any]:
    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)

    payload = json.loads(presentation_json)
    slides = payload.get("slides") or []
    if not isinstance(slides, list) or not slides:
        raise ValueError("演示文稿至少需要一页幻灯片")

    deck_title = str(payload.get("title") or "演示文稿").strip()
    deck_subtitle = str(payload.get("subtitle") or "").strip()

    presentation = Presentation()
    presentation.slide_width = PAGE_WIDTH
    presentation.slide_height = PAGE_HEIGHT

    for index, slide_payload in enumerate(slides):
        slide_title = str(slide_payload.get("title") or f"第 {index + 1} 页").strip()
        bullets_raw = slide_payload.get("bullets") or []
        if not isinstance(bullets_raw, list):
            bullets_raw = []
        bullets = [str(item).strip() for item in bullets_raw if str(item).strip()][:5]

        if index == 0:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            build_title_slide(slide, deck_title or slide_title, deck_subtitle, bullets)
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        build_content_slide(slide, slide_title, bullets, index + 1)

    presentation.save(target)
    return {
        "ok": True,
        "outputPath": str(target),
        "slideCount": len(presentation.slides),
    }


def build_title_slide(slide, title: str, subtitle: str, bullets: list[str]) -> None:
    apply_background(slide, COLOR_BG)
    add_corner_glow(slide)
    add_top_bar(slide, 0.55, COLOR_ACCENT)

    badge = slide.shapes.add_textbox(Inches(0.78), Inches(0.72), Inches(2.2), Inches(0.4))
    badge_tf = badge.text_frame
    badge_tf.text = "KNOWLEDGEOS PRESENTATION"
    style_paragraph(badge_tf.paragraphs[0], 11, COLOR_MUTED, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.45), Inches(7.6), Inches(1.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.text = title
    style_paragraph(title_tf.paragraphs[0], 28, COLOR_TEXT, bold=True)

    subtitle_box = slide.shapes.add_textbox(Inches(0.82), Inches(3.1), Inches(6.8), Inches(1.0))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.text = subtitle or "基于资料自动整理生成"
    style_paragraph(subtitle_tf.paragraphs[0], 16, COLOR_MUTED)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.85),
        Inches(1.35),
        Inches(3.55),
        Inches(4.8),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_PANEL
    panel.line.color.rgb = COLOR_ACCENT_SOFT
    panel.line.width = Pt(1.2)

    panel_title = slide.shapes.add_textbox(Inches(9.18), Inches(1.72), Inches(2.8), Inches(0.45))
    panel_title_tf = panel_title.text_frame
    panel_title_tf.text = "本次内容"
    style_paragraph(panel_title_tf.paragraphs[0], 15, COLOR_TEXT, bold=True)

    panel_list = slide.shapes.add_textbox(Inches(9.18), Inches(2.25), Inches(2.7), Inches(3.4))
    panel_tf = panel_list.text_frame
    panel_tf.word_wrap = True
    if bullets:
        panel_tf.text = bullets[0]
        style_paragraph(panel_tf.paragraphs[0], 18, COLOR_TEXT, bold=True)
        for bullet in bullets[1:4]:
            paragraph = panel_tf.add_paragraph()
            paragraph.text = bullet
            style_paragraph(paragraph, 14, COLOR_MUTED)
            paragraph.level = 0
    else:
        panel_tf.text = "项目背景\n核心要点\n结论总结"
        style_paragraph(panel_tf.paragraphs[0], 18, COLOR_TEXT, bold=True)
        for item in ["核心要点", "结论总结"]:
            paragraph = panel_tf.add_paragraph()
            paragraph.text = item
            style_paragraph(paragraph, 14, COLOR_MUTED)


def build_content_slide(slide, title: str, bullets: list[str], page_number: int) -> None:
    apply_background(slide, COLOR_BG_SOFT)
    add_top_bar(slide, 0.18, COLOR_ACCENT_SOFT)
    add_side_block(slide)

    title_box = slide.shapes.add_textbox(Inches(0.92), Inches(0.9), Inches(8.2), Inches(1.15))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.text = title
    style_paragraph(title_tf.paragraphs[0], 26, COLOR_TEXT, bold=True)

    index_box = slide.shapes.add_textbox(Inches(11.55), Inches(0.88), Inches(1.0), Inches(0.4))
    index_tf = index_box.text_frame
    index_tf.text = f"{page_number:02d}"
    style_paragraph(index_tf.paragraphs[0], 14, COLOR_MUTED, bold=True, align="right")

    content_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.88),
        Inches(1.72),
        Inches(10.25),
        Inches(4.95),
    )
    content_panel.fill.solid()
    content_panel.fill.fore_color.rgb = RGBColor(25, 34, 49)
    content_panel.line.color.rgb = RGBColor(45, 58, 82)
    content_panel.line.width = Pt(1.0)

    bullet_box = slide.shapes.add_textbox(Inches(1.28), Inches(2.1), Inches(9.25), Inches(4.1))
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True

    if bullets:
        bullet_tf.text = bullets[0]
        style_paragraph(bullet_tf.paragraphs[0], 22, COLOR_TEXT, bold=True)
        for bullet in bullets[1:]:
            paragraph = bullet_tf.add_paragraph()
            paragraph.text = bullet
            style_paragraph(paragraph, 20, COLOR_TEXT)
            paragraph.level = 0
            paragraph.space_before = Pt(10)
    else:
        bullet_tf.text = "暂无可展示要点"
        style_paragraph(bullet_tf.paragraphs[0], 20, COLOR_MUTED)

    footer = slide.shapes.add_textbox(Inches(0.96), Inches(6.88), Inches(4.6), Inches(0.35))
    footer_tf = footer.text_frame
    footer_tf.text = "KnowledgeOS 自动生成"
    style_paragraph(footer_tf.paragraphs[0], 10, COLOR_MUTED)


def apply_background(slide, color: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color


def add_corner_glow(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(9.9),
        Inches(-0.65),
        Inches(4.2),
        Inches(4.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(45, 82, 150)
    shape.fill.transparency = 0.28
    shape.line.fill.background()


def add_top_bar(slide, top: float, color: RGBColor) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(top),
        PAGE_WIDTH,
        Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_side_block(slide) -> None:
    side = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(11.45),
        Inches(1.72),
        Inches(1.05),
        Inches(4.95),
    )
    side.fill.solid()
    side.fill.fore_color.rgb = RGBColor(33, 48, 72)
    side.line.fill.background()


def style_paragraph(paragraph, size: int, color: RGBColor, bold: bool = False, align: str | None = None) -> None:
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    if align == "right":
        from pptx.enum.text import PP_ALIGN

        paragraph.alignment = PP_ALIGN.RIGHT
