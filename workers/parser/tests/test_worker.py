from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfWriter


ROOT = Path(__file__).resolve().parents[3]
WORKER = ROOT / "workers" / "parser" / "main.py"


def test_health() -> None:
    completed = subprocess.run(
        [sys.executable, str(WORKER), "health"],
        capture_output=True,
        check=True,
        text=True,
    )
    payload = json.loads(completed.stdout)
    assert payload["ok"] is True


def test_parse_file(tmp_path: Path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# 示例\n\n内容", encoding="utf-8")

    completed = subprocess.run(
        [
            sys.executable,
            str(WORKER),
            "parse_file",
            "--file-path",
            str(source),
            "--source-type",
            "md",
        ],
        capture_output=True,
        check=True,
        text=True,
    )
    payload = json.loads(completed.stdout)
    assert payload["ok"] is True
    assert payload["manifest"]["title"] == "sample"


def test_parse_docx(tmp_path: Path) -> None:
    source = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("文档标题", level=1)
    document.add_paragraph("这里是正文。")
    document.save(source)

    completed = subprocess.run(
        [
            sys.executable,
            str(WORKER),
            "parse_file",
            "--file-path",
            str(source),
            "--source-type",
            "docx",
        ],
        capture_output=True,
        check=True,
        text=True,
    )
    payload = json.loads(completed.stdout)
    assert payload["ok"] is True
    assert "文档标题" in payload["markdown"]


def test_parse_pptx(tmp_path: Path) -> None:
    source = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "第一页"
    slide.placeholders[1].text = "要点一\n要点二"
    presentation.save(source)

    completed = subprocess.run(
        [
            sys.executable,
            str(WORKER),
            "parse_file",
            "--file-path",
            str(source),
            "--source-type",
            "pptx",
        ],
        capture_output=True,
        check=True,
        text=True,
    )
    payload = json.loads(completed.stdout)
    assert payload["ok"] is True
    assert "第一页" in payload["markdown"]


def test_generate_pptx(tmp_path: Path) -> None:
    output = tmp_path / "generated.pptx"
    payload = {
        "title": "测试演示文稿",
        "subtitle": "副标题",
        "slides": [
            {"title": "测试演示文稿", "bullets": []},
            {"title": "核心结论", "bullets": ["要点一", "要点二", "要点三"]},
        ],
    }

    completed = subprocess.run(
        [
            sys.executable,
            str(WORKER),
            "generate_pptx",
            "--output-path",
            str(output),
            "--presentation-json",
            json.dumps(payload, ensure_ascii=False),
        ],
        capture_output=True,
        check=True,
        text=True,
    )
    result = json.loads(completed.stdout)
    assert result["ok"] is True
    assert result["slideCount"] == 2
    assert output.exists()

    presentation = Presentation(output)
    assert len(presentation.slides) == 2
    first_slide_text = "\n".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text"))
    assert "测试演示文稿" in first_slide_text


def test_parse_pdf(tmp_path: Path) -> None:
    source = tmp_path / "sample.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with source.open("wb") as handle:
        writer.write(handle)

    completed = subprocess.run(
        [
            sys.executable,
            str(WORKER),
            "parse_file",
            "--file-path",
            str(source),
            "--source-type",
            "pdf",
        ],
        capture_output=True,
        check=True,
        text=True,
    )
    payload = json.loads(completed.stdout)
    assert payload["ok"] is True
    assert payload["manifest"]["sourceType"] == "pdf"
    assert "当前 PDF 没有提取到可转换的正文内容" in payload["markdown"]
    assert payload["manifest"]["warnings"]
